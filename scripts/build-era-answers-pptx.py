#!/usr/bin/env python3
"""Build the 9-slide ERA reusability deck as a real .pptx — one slide per question.

Run:  python3 scripts/build-era-answers-pptx.py   (needs python-pptx)
Source text: interop-europe/answers.md  ·  page: scripts/assets/era-interop-answers.html

Styling follows the ERA deck template: navy title/closing panels, white content
slides with a numbered navy chip and the agency mark, a light-weight grey slide
title set right, and navy-headed tables on pale blue rows.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import pathlib

NAVY   = RGBColor(0x2E, 0x43, 0x72)
CHIP   = RGBColor(0x2F, 0x4B, 0x7C)
FILL   = RGBColor(0xDC, 0xE6, 0xF1)
FILL2  = RGBColor(0xED, 0xF2, 0xF9)
HEAD   = RGBColor(0x59, 0x59, 0x59)
INK    = RGBColor(0x26, 0x26, 0x26)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
AMBER  = RGBColor(0x8A, 0x56, 0x10)
AMBBG  = RGBColor(0xFB, 0xF1, 0xE0)
GREEN  = RGBColor(0x25, 0x5C, 0x45)
GRNBG  = RGBColor(0xE8, 0xF1, 0xEC)
MUTED  = RGBColor(0x5B, 0x65, 0x77)
FONT   = "Calibri"
MONO   = "Consolas"

LOGO = pathlib.Path(__file__).resolve().parent / "assets" / "era-logo.png"
W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def box(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.text_frame.clear()
    return sh


def text(slide, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font=FONT, space=6, line=1.25):
    """runs: str, or list of paragraphs; a paragraph is a str or list of (txt, **fmt)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.line_spacing = line
        parts = para if isinstance(para, list) else [(para, {})]
        for txt, fmt in parts:
            r = p.add_run(); r.text = txt
            f = r.font
            f.name = fmt.get("font", font)
            f.size = Pt(fmt.get("size", size))
            f.bold = fmt.get("bold", bold)
            f.italic = fmt.get("italic", False)
            f.color.rgb = fmt.get("color", color)
    return tb


def header(slide, num, title):
    """Navy numbered chip + agency mark + right-set light title."""
    box(slide, Inches(.62), Inches(.42), Inches(.5), Inches(.5), fill=CHIP)
    text(slide, Inches(.62), Inches(.53), Inches(.5), Inches(.4), str(num),
         size=15, color=WHITE, align=PP_ALIGN.CENTER)
    slide.shapes.add_picture(str(LOGO), Inches(1.30), Inches(.36), height=Inches(.62))
    text(slide, Inches(5.4), Inches(.36), Inches(7.3), Inches(1.0), title,
         size=27, color=HEAD, align=PP_ALIGN.RIGHT, line=1.06)


def footer(slide, left, num):
    text(slide, Inches(.62), Inches(6.98), Inches(9.5), Inches(.3), left,
         size=9.5, color=RGBColor(0x8A, 0x93, 0xA6))
    text(slide, Inches(12.0), Inches(6.98), Inches(.7), Inches(.3), str(num),
         size=9.5, color=RGBColor(0x8A, 0x93, 0xA6), align=PP_ALIGN.RIGHT)


def callout(slide, x, y, w, h, kind, parts):
    bg, bar, hue = {"amber": (AMBBG, AMBER, AMBER),
                    "green": (GRNBG, GREEN, GREEN),
                    "navy":  (FILL, NAVY, NAVY)}[kind]
    box(slide, x, y, w, h, fill=bg)
    box(slide, x, y, Inches(.05), h, fill=bar)
    runs = [(t, dict(f, color=hue) if f.get("bold") else f) for t, f in parts]
    text(slide, x + Inches(.22), y + Inches(.14), w - Inches(.42), h - Inches(.2),
         [runs], size=13.5, line=1.22)


def table(slide, x, y, w, rows, widths, size=12.5, rh=Inches(.32)):
    n, m = len(rows), len(rows[0])
    shp = slide.shapes.add_table(n, m, x, y, w, rh * n)
    tbl = shp.table
    tbl.first_row = True
    for j, fr in enumerate(widths):
        tbl.columns[j].width = Emu(int(w * fr))
    for i, row in enumerate(rows):
        tbl.rows[i].height = rh
        for j, cell in enumerate(row):
            c = tbl.cell(i, j)
            c.text = ""
            c.margin_left = c.margin_right = Inches(.09)
            c.margin_top = c.margin_bottom = Inches(.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = NAVY if i == 0 else (FILL if i % 2 else FILL2)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if (j and str(cell).replace(",", "").replace("%", "").isdigit()) else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(cell)
            r.font.size = Pt(size)
            r.font.name = FONT
            r.font.bold = (i == 0)
            r.font.color.rgb = WHITE if i == 0 else INK
    return shp



# ── extra visual helpers ──────────────────────────────────────────
ACCENT = RGBColor(0x3F, 0x7B, 0xD1)
TEAL   = RGBColor(0x1F, 0x6B, 0x4A)


def tiles(slide, x, y, w, items, h=Inches(1.12), gap=Inches(.16), size=30, lab=10.5):
    """A row of big-number stat tiles: [(value, label), …]."""
    n = len(items)
    tw = (w - gap * (n - 1)) / n
    for i, (val, label) in enumerate(items):
        tx = x + (tw + gap) * i
        box(slide, tx, y, tw, h, fill=FILL2)
        box(slide, tx, y, tw, Inches(.045), fill=ACCENT)
        text(slide, tx + Inches(.14), y + Inches(.19), tw - Inches(.28), Inches(.5),
             val, size=size, bold=True, color=NAVY, align=PP_ALIGN.CENTER, line=1.0)
        text(slide, tx + Inches(.1), y + h - Inches(.42), tw - Inches(.2), Inches(.36),
             label, size=lab, color=MUTED, align=PP_ALIGN.CENTER, line=1.12)


def bars(slide, x, y, w, rows, bar_h=Inches(.26), gap=Inches(.13),
         label_w=Inches(2.05), value_w=Inches(1.05), hue=ACCENT, note=None):
    """Horizontal bar chart. rows: [(label, value, display)] — value drives length."""
    top = max(v for _, v, _ in rows) or 1
    track = w - label_w - value_w
    for i, (label, val, disp) in enumerate(rows):
        by = y + (bar_h + gap) * i
        text(slide, x, by - Inches(.02), label_w - Inches(.12), bar_h,
             label, size=11.5, color=INK, align=PP_ALIGN.RIGHT, line=1.0)
        box(slide, x + label_w, by, track, bar_h, fill=FILL2)
        bw = max(int(track * val / top), Emu(9000))
        box(slide, x + label_w, by, Emu(int(bw)), bar_h, fill=hue if i == 0 else FILL)
        text(slide, x + label_w + track + Inches(.1), by - Inches(.02), value_w, bar_h,
             disp, size=11.5, bold=(i == 0), color=NAVY if i == 0 else INK, line=1.0)
    if note:
        text(slide, x + label_w, y + (bar_h + gap) * len(rows) + Inches(.06),
             w - label_w, Inches(.3), [[(note, {"italic": True})]], size=10, color=MUTED)


def bullets(slide, x, y, w, h, items, size=13.5, space=9, lead=None):
    """items: [(lead, rest)] — lead set navy bold, rest normal."""
    paras = []
    for it in items:
        if isinstance(it, tuple):
            paras.append([("▪  ", {"color": ACCENT, "bold": True}),
                          (it[0], {"bold": True, "color": NAVY}), (it[1], {})])
        else:
            paras.append([("▪  ", {"color": ACCENT, "bold": True}), (it, {})])
    text(slide, x, y, w, h, paras, size=size, space=space, line=1.28)


def slide(num, title, sub=None):
    sl = prs.slides.add_slide(BLANK)
    header(sl, num, title)
    if sub:
        text(sl, Inches(.62), Inches(1.32), Inches(12.1), Inches(.55), sub,
             size=15.5, color=NAVY, bold=True, line=1.2)
    footer(sl, "ERA Ontology v3.3.4 · EUPL 1.2 · figures measured 22 August 2026", num + 1)
    return sl


# ── 1 ─────────────────────────────────────────────────────── title
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, fill=NAVY)
text(s, Inches(.95), Inches(1.95), Inches(9.6), Inches(1.05),
     [[("Reusing the ", {"size": 46}), ("ERA Ontology", {"size": 46, "bold": True})]],
     color=WHITE, line=1.02, space=2)
text(s, Inches(.95), Inches(3.05), Inches(9.6), Inches(.6),
     "Seven answers for the Interoperable Europe assessment", size=21, color=WHITE)
box(s, Inches(.95), Inches(3.9), Inches(4.2), Pt(1), fill=RGBColor(0x8E, 0xA2, 0xC6))
text(s, Inches(.95), Inches(4.15), Inches(10.2), Inches(1.2),
     ["Every figure is a legal reference, a measurement of the published artefacts "
      "(ontology v3.3.4, era-shapes, era-skos), or a measurement against the live "
      "knowledge graph on 22 August 2026.",
      "Parameter counts exclude owl:deprecated terms. Where the artefacts and the "
      "deployed graph disagree, both numbers are given."],
     size=13, color=RGBColor(0xD3, 0xDD, 0xEE), line=1.3)
prs.slides[0].shapes.add_picture(str(LOGO), Inches(.95), Inches(.75), height=Inches(.72))

# ── 2 ─────────────────────────────────── Q1 reusability / evidence
s = slide(1, "Would it make sense for public authorities\nto look at reusability?",
          "Yes — because the ontology is not an interpretation of the law. It is the law in machine-readable form.")
callout(s, Inches(.62), Inches(2.0), Inches(12.1), Inches(.62), "navy",
        [("RINF parameter 1.1.1.1.4.1 “Nominal track gauge” ", {"bold": True}),
         ("is ", {}), ("era:wheelSetGauge", {"font": MONO}),
         (". 292 live properties carry their legal index, so the trace from legal text to "
          "data field is machine-checkable, not a matter of documentation discipline.", {})])
tiles(s, Inches(.62), Inches(2.86), Inches(12.1),
      [("47.4 M", "triples in the\nlive graph"), ("27", "countries\npublishing"),
       ("54", "national\ndatasets"), ("725,637", "running\ntracks"),
       ("18,279", "SKOS concepts\nin 423 schemes"), ("1,029", "SHACL shapes\n+ 323 constraints")])
text(s, Inches(.62), Inches(4.28), Inches(5.6), Inches(.3),
     "The same design across four legal registers", size=13, bold=True, color=NAVY)
bars(s, Inches(.62), Inches(4.72), Inches(5.6),
     [("era:rinfIndex", 292, "292"), ("era:eratvIndex", 174, "174"),
      ("tsiOPEAppendixD2Index", 67, "67")],
     label_w=Inches(1.95), value_w=Inches(.6),
     note="live properties carrying a legal index (RINF · ERATV · Route Book)")
callout(s, Inches(6.65), Inches(4.28), Inches(6.07), Inches(1.62), "green",
        [("The chain is followable end to end. ", {"bold": True}),
         ("era-lex holds 208,517 triples over 7,317 legal acts in ELI, and 420 ontology "
          "terms cite 79 acts through dcterms:source. A reuser goes from data field → "
          "numbered parameter → provision in the Official Journal with no human in the "
          "loop. This is the part most worth copying, and it is domain-neutral.", {})])

# ── 3 ─────────────────────────────────── Q2 challenge and lessons
s = slide(2, "What was the main challenge,\nand what were the lessons?",
          "The vocabulary was the easy part. The challenge is that a shared vocabulary does not "
          "guarantee shared practice — and the gap is invisible until you measure it.")
text(s, Inches(.62), Inches(2.16), Inches(6.5), Inches(.3),
     "Retired parameters still carrying live data, by publisher", size=13, bold=True, color=NAVY)
bars(s, Inches(.62), Inches(2.6), Inches(6.5),
     [("Germany", 156064, "156,064"), ("Switzerland", 63159, "63,159"),
      ("Czechia", 6538, "6,538"), ("Netherlands", 3409, "3,409"),
      ("Sweden", 3095, "3,095"), ("Spain", 2607, "2,607"),
      ("Hungary", 2250, "2,250"), ("Ireland", 1670, "1,670"), ("France", 592, "592")],
     label_w=Inches(1.5), value_w=Inches(.95),
     note="9 countries · 17 retired parameters · 237,781 statements in total")
bullets(s, Inches(7.5), Inches(2.16), Inches(5.22), Inches(3.4), [
    ("Divergence is silent. ", "Croatia publishes op-types under a parallel concept "
     "scheme with typos (“Tehnical change”); the SHACL check passes because the value "
     "shape tests membership of a scheme, not of the canonical one."),
    ("Deprecation is a migration, ", "not an annotation. 37 RINF parameters are retired "
     "and 9 countries still publish them."),
    ("An identifier only counts if it resolves. ", "Belgium minted 0088 for INFRABEL "
     "where the register already had 1976 — name, VAT and 7 roles on the real one, "
     "none on the copy."),
    ("Validate the vocabulary, not only the data, ", "and resolve value sets from a "
     "trusted source rather than from the graph under test."),
])
callout(s, Inches(.62), Inches(6.06), Inches(12.1), Inches(.72), "green",
        [("These findings are now enforced, not just recorded. ", {"bold": True}),
         ("All three tools built on this graph exclude owl:deprecated properties, and the "
          "nightly catalogue rebuild reports which datasets still carry retired data — so "
          "the gap stays visible instead of decaying back into folklore.", {})])

# ── 4 ─────────────────────────────────── Q3 collaboration
s = slide(3, "Do you envisage collaboration with other\nMember States or Union entities?",
          "Collaboration is already the operating model, not an aspiration.")
bullets(s, Inches(.62), Inches(2.15), Inches(7.3), Inches(4.2), [
    ("Standardisation bodies. ", "Memorandum of Intent with railML.org signed 30 May 2023; "
     "ontologies jointly published and railML→RINF transformations developed."),
    ("The Publications Office. ", "ERA mints no country codes — the graph points straight at "
     "the EU authority tables. Reusing existing EU semantic assets is a deliberate choice."),
    ("W3C and OGC. ", "GeoSPARQL, SKOS, PROV, ORG and OWL-Time are imported rather than "
     "geometry, provenance and organisations being reinvented."),
    ("Member States and infrastructure managers, ", "through the National Registration "
     "Entities the RINF Regulation requires, and a public GitLab tracker where every "
     "change request is visible."),
])
callout(s, Inches(8.15), Inches(2.15), Inches(4.57), Inches(1.62), "green",
        [("Adoption beyond the mandate. ", {"bold": True}),
         ("Bane NOR, the Norwegian infrastructure manager, publishes its own station data "
          "with the ERA ontology on its own Linked Data server — outside the EU compliance "
          "perimeter, visible in the shared graph. Voluntary adoption is the clearest "
          "evidence that the model is reusable on its merits.", {})])
text(s, Inches(.62), Inches(4.55), Inches(12.1), Inches(.3),
     "Next steps we would welcome partners on", size=13, bold=True, color=NAVY)
text(s, Inches(.62), Inches(4.95), Inches(12.1), Inches(1.0),
     ["European Mobility Data Space and National Access Points  ·  publishing ERATV, EVR and "
      "ERADIS data as openly as RINF already is  ·  closing the value-set validation gaps  ·  "
      "contributing the pattern back to SEMIC as a reusable design rather than a rail artefact."],
     size=13, line=1.3)

# ── 5 ─────────────────────────────────── Q4 who can use it
s = slide(4, "Is it suitable for local, regional\nor national administrations?",
          "Yes, at every level — because the unit of publication is the dataset, not the country.")
text(s, Inches(.62), Inches(2.2), Inches(6.9), Inches(.3),
     "Publishing organisations per country — 54 datasets, 27 countries", size=13, bold=True, color=NAVY)
bars(s, Inches(.62), Inches(2.66), Inches(6.9),
     [("Italy", 9, "9"), ("Austria", 8, "8"), ("France", 7, "7"), ("Sweden", 4, "4"),
      ("Germany", 2, "2"), ("Finland", 2, "2"), ("Switzerland", 2, "2"),
      ("20 other countries", 1, "1 each")],
     label_w=Inches(1.85), value_w=Inches(.85), bar_h=Inches(.28),
     note="small regional and private managers publish alongside national incumbents")
bullets(s, Inches(7.85), Inches(2.2), Inches(4.87), Inches(3.6), [
    ("Not a national-champions model. ", "No entity has to be large to participate; every "
     "publisher appears in the same cross-border queries."),
    ("The granularity supports local asset management, ", "not only strategic planning — "
     "platforms, sidings, level crossings, tunnels, bridges, signals and kilometric posts "
     "are all modelled."),
    ("Value arrives before anyone else joins. ", "A regional authority adopting it for its "
     "own network gets consistent asset descriptions validated against the shipped SHACL "
     "shapes on day one."),
])

# ── 6 ─────────────────────────────────── Q5 governance
s = slide(5, "Does it need central governance?",
          "Two things must be separated, and the honest answer differs for each.")
box(s, Inches(.62), Inches(2.35), Inches(5.95), Inches(2.05), fill=FILL2)
box(s, Inches(.62), Inches(2.35), Inches(5.95), Inches(.05), fill=NAVY)
text(s, Inches(.92), Inches(2.60), Inches(5.35), Inches(.4),
     "The vocabulary — central, non-negotiable", size=15, bold=True, color=NAVY)
text(s, Inches(.92), Inches(3.08), Inches(5.35), Inches(1.25),
     ["A shared meaning must have exactly one owner: one version line (v3.3.4), one licence "
      "(EUPL 1.2), a public repository and an archival DOI.",
      "Without that you do not get one vocabulary used 27 times. You get 27 dialects.",
      "And it has to extend to what is loaded, not stop at what is published — the governed "
      "artefact is clean, yet the deployed graph carries concepts that were never in it."],
     size=12.5, line=1.3, space=7)
box(s, Inches(6.77), Inches(2.35), Inches(5.95), Inches(2.05), fill=GRNBG)
box(s, Inches(6.77), Inches(2.35), Inches(5.95), Inches(.05), fill=GREEN)
text(s, Inches(7.07), Inches(2.60), Inches(5.35), Inches(.4),
     "The deployment — decentralised, immediate", size=15, bold=True, color=GREEN)
text(s, Inches(7.07), Inches(3.08), Inches(5.35), Inches(1.25),
     ["Benefit accrues per publisher. Each entity that adopts the vocabulary gets value "
      "before anyone else joins: its own data becomes queryable, comparable year-on-year "
      "and testable against the SHACL shapes.",
      "Bane NOR demonstrates exactly this — it derives value publishing essentially alone, "
      "outside the EU mandate."],
     size=12.5, line=1.3, space=7)
callout(s, Inches(.62), Inches(4.72), Inches(12.1), Inches(.8), "amber",
        [("The realistic caveat. ", {"bold": True}),
         ("Cross-border value requires both critical mass and conformance. Partial deployment "
          "with local variation is the worst of both worlds — it produces the appearance of "
          "interoperability without the substance, because queries return plausible but "
          "incomplete answers. Partial deployment with conformance is genuinely useful from "
          "the first publisher onward.", {})])

# ── 7 ─────────────────────────────────── Q6 cross-border + sovereignty
s = slide(6, "Cross-border data exchange\nand EU sovereignty",
          "Route compatibility spans three registers at once — so one shared vocabulary turns an "
          "integration project into a query.")
text(s, Inches(.62), Inches(2.3), Inches(6.4), Inches(.3),
     "One question, the whole European gauge landscape — RINF 1.1.1.1.4.1",
     size=13, bold=True, color=NAVY)
bars(s, Inches(.62), Inches(2.76), Inches(6.4),
     [("1435  standard", 697901, "697,901"), ("1668  Iberian", 13275, "13,275"),
      ("1524  Finnish/Baltic", 2813, "2,813"), ("1000  metre", 2746, "2,746"),
      ("1520  ex-Soviet", 1826, "1,826"), ("1600  Irish", 274, "274"),
      ("760 / 750", 135, "135")],
     label_w=Inches(1.95), value_w=Inches(.95),
     note="running tracks per gauge across 27 countries, from one query")
bullets(s, Inches(7.35), Inches(2.3), Inches(5.37), Inches(3.9), [
    ("Custody stays national. ", "Each Member State publishes into its own dataset. What is "
     "shared is the meaning, not the ownership — there is no central database taking "
     "possession of national data."),
    ("EU-controlled identifiers. ", "Whoever controls the identifiers controls who can join "
     "the data. Here that authority is European, the URIs are persistent and publicly "
     "governed — the most durable form of digital sovereignty on offer."),
    ("No vendor lock-in. ", "RDF, SPARQL, SKOS, SHACL, GeoSPARQL — open standards, multiple "
     "implementations, EUPL 1.2, archived on Zenodo."),
    ("Open access enables scrutiny. ", "Anyone can verify a claim about the register "
     "directly. The data-quality findings in §2 were found exactly that way."),
])
callout(s, Inches(.62), Inches(5.98), Inches(6.4), Inches(.92), "navy",
        [("Borders are first-class objects: ", {"bold": True}),
         ("288 reference border points, 427 border and 433 domestic-border operational "
          "points — and published queries expose 2,244 disconnected points in ERA’s own "
          "register.", {})])

# ── 8 ─────────────────────────────────── Q7 first step
s = slide(7, "What is the first step\nto reuse the solution?",
          "Query it before committing to anything. No account, no licence negotiation, no data pipeline.")
box(s, Inches(.62), Inches(2.12), Inches(6.4), Inches(1.62), fill=RGBColor(0xF3, 0xF5, 0xF9))
box(s, Inches(.62), Inches(2.12), Inches(.05), Inches(1.62), fill=ACCENT)
text(s, Inches(.86), Inches(2.28), Inches(6.05), Inches(1.35),
     ["curl -sS -X POST \\",
      "  \"https://graph.data.era.europa.eu/repositories/rinf-plus\" \\",
      "  -H \"Content-Type: application/sparql-query\" \\",
      "  --data-binary 'SELECT ?value (COUNT(?track) AS ?n) WHERE {",
      "    ?p era:rinfIndex \"1.1.1.1.4.1\" . ?track ?p ?value }",
      "    GROUP BY ?value ORDER BY DESC(?n)'"],
     size=10.5, font=MONO, color=RGBColor(0x1F, 0x33, 0x55), line=1.24, space=1)
text(s, Inches(.62), Inches(3.88), Inches(6.4), Inches(.3),
     "Runs as-is and returns the gauge table on the previous slide.", size=11.5, color=MUTED)
text(s, Inches(7.25), Inches(2.12), Inches(5.47), Inches(.3),
     "Then, in order", size=13, bold=True, color=NAVY)
steps = [("Run the published queries first — ", "38 in the catalogue, 47 in three notebooks."),
         ("Read the vocabulary ", "and the application guide that maps each parameter."),
         ("Locate your own parameters ", "by rinfIndex; what does not map is national, or a gap."),
         ("Adopt the identifiers before modelling. ", "Cheapest interoperability there is."),
         ("Adopt the code lists before the classes ", "— and pin the version you validated."),
         ("Validate early with the shipped SHACL shapes, ", "in your own pipeline."),
         ("Publish into your own dataset, ", "keeping custody."),
         ("Engage the maintainers ", "through the public issue tracker.")]
paras = []
for i, (lead, rest) in enumerate(steps, 1):
    paras.append([(f"{i}  ", {"bold": True, "color": ACCENT}),
                  (lead, {"bold": True, "color": NAVY}), (rest, {})])
text(s, Inches(7.25), Inches(2.55), Inches(5.47), Inches(4.0), paras,
     size=12.5, space=7, line=1.24)
callout(s, Inches(.62), Inches(4.42), Inches(6.4), Inches(1.55), "green",
        [("If your domain is not rail, ", {"bold": True}),
         ("the reusable asset is the design, not the classes: a legal annex expressed as "
          "indexed properties, governed code lists shipped as SKOS, validation shapes "
          "shipped with the vocabulary, and open query access. That pattern is what we "
          "would encourage assessors to consider transferable — to energy metering, "
          "building permits, or any register a legal act mandates across 27 jurisdictions.", {})])

# ── 9 ─────────────────────────────────────────────── sources
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, fill=NAVY)
s.shapes.add_picture(str(LOGO), Inches(.95), Inches(.72), height=Inches(.66))
text(s, Inches(.95), Inches(1.85), Inches(11.4), Inches(.7),
     "Sources — every figure above is traceable", size=30, bold=True, color=WHITE)
box(s, Inches(.95), Inches(2.72), Inches(4.2), Pt(1), fill=RGBColor(0x8E, 0xA2, 0xC6))
text(s, Inches(.95), Inches(3.02), Inches(5.6), Inches(2.6),
     ["ERA Ontology v3.3.4 — rinf.data.era.europa.eu/era-vocabulary",
      "Application guides — RINF, ERATV, EVR, ERADIS (same host)",
      "Published artefacts — era-shapes, era-skos, era-telem-skos",
      "Live knowledge graph — graph.data.era.europa.eu",
      "Data Stories — 38 queries + 3 notebooks",
      "DOI 10.5281/zenodo.15089005 · EUPL 1.2"],
     size=12, color=RGBColor(0xD3, 0xDD, 0xEE), line=1.34, space=5)
text(s, Inches(6.85), Inches(3.02), Inches(5.6), Inches(2.6),
     ["Reg. (EU) 2019/777, amended by (EU) 2023/1694 — RINF",
      "Decision 2011/665/EU; Reg. (EU) 2019/776 — ERATV",
      "Reg. (EU) 2023/1695; Decision (EU) 2018/1614 — EVR",
      "Dir. (EU) 2016/797 and 2016/798; Reg. (EU) 2016/796 — ERADIS",
      "Interoperable Europe Portal — ERA Vocabulary solution",
      "Logo: ERA-rgb-300dpi.jpg, Wikimedia Commons, public domain, by ERA"],
     size=12, color=RGBColor(0xD3, 0xDD, 0xEE), line=1.34, space=5)
text(s, Inches(.95), Inches(6.1), Inches(11.4), Inches(.5),
     "Full answers: gatemezing.github.io/onto-dissemination/interopable-eu-portal-answers.html",
     size=13, color=WHITE)

out = (pathlib.Path(__file__).resolve().parent.parent /
       "interop-europe" / "ERA-ontology-reusability.pptx")
prs.save(out)
print(f"saved {out.name}: {out.stat().st_size:,} bytes, {len(prs.slides._sldIdLst)} slides")
